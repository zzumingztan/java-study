"""Extract text from all chapter PPT/PDF files and save as JSON."""
import json
import os
import sys

BASE_DIR = r"C:\Users\21268\Desktop\java实验"
OUTPUT = os.path.join(os.path.dirname(__file__), "chapters_text.json")

# Map files to chapter numbers
FILES = {
    1: "第1章 Java 语言基础知识.pptx",
    2: "第2章 类与对象的基本概念.pptx",
    3: "第3章 类的重用(1).pptx",
    4: "第4章接口与多态.pptx",
    5: "第5章异常处理与输入输出流-已经完成.pptx",
    6: "第6章-集合框架.pdf",
    7: "第7章 图形用户界面.pdf",
    8: "第8章 多线程编程.pdf",
}

CHAPTER_TITLES = {
    1: "Java 语言基础知识",
    2: "类与对象的基本概念",
    3: "类的重用",
    4: "接口与多态",
    5: "异常处理与输入输出流",
    6: "集合框架",
    7: "图形用户界面",
    8: "多线程编程",
}


def extract_pptx(filepath):
    """Extract text from a PPTX file."""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    lines.append(" | ".join(row_text))
        if lines:
            slides_text.append({"slide": i, "content": "\n".join(lines)})
    return slides_text


def extract_pdf(filepath):
    """Extract text from a PDF file."""
    from PyPDF2 import PdfReader
    reader = PdfReader(filepath)
    pages_text = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages_text.append({"page": i, "content": text.strip()})
    return pages_text


def main():
    chapters = {}
    for num, filename in sorted(FILES.items()):
        filepath = os.path.join(BASE_DIR, filename)
        if not os.path.exists(filepath):
            print(f"  [WARN] Chapter {num}: file not found: {filepath}")
            continue

        ext = os.path.splitext(filename)[1].lower()
        print(f"Extracting Chapter {num}: {filename} ...", end=" ")
        try:
            if ext == ".pptx":
                slides = extract_pptx(filepath)
            elif ext == ".pdf":
                slides = extract_pdf(filepath)
            else:
                print(f"SKIP (unknown format)")
                continue

            full_text = "\n\n".join(s["content"] for s in slides)
            chapters[str(num)] = {
                "title": CHAPTER_TITLES[num],
                "filename": filename,
                "slides_count": len(slides),
                "slides": slides,
                "full_text": full_text,
                "char_count": len(full_text),
            }
            print(f"OK ({len(slides)} slides/pages, {len(full_text)} chars)")
        except Exception as e:
            print(f"FAIL: {e}")

    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    with open(OUTPUT, "w", encoding="utf-8") as f:
        json.dump(chapters, f, ensure_ascii=False, indent=2)
    print(f"\nSaved to: {OUTPUT}")
    for num, ch in sorted(chapters.items(), key=lambda x: int(x[0])):
        print(f"  Ch{num}: {ch['title']} — {ch['char_count']} chars")


if __name__ == "__main__":
    main()
